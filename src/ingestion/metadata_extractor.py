"""
Enhanced metadata extraction for files
"""
from typing import Dict, Any, Optional, List, Set
from pathlib import Path
from datetime import datetime
import hashlib
import json
import re
from dataclasses import dataclass, asdict
import pandas as pd

try:
    from openpyxl import load_workbook
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False


@dataclass
class FileMetadata:
    """Rich metadata for a file"""
    file_id: str
    file_path: str
    file_name: str
    file_type: str
    file_size: int
    extension: str
    created_at: Optional[str] = None
    modified_at: Optional[str] = None
    author: Optional[str] = None
    title: Optional[str] = None
    subject: Optional[str] = None
    content_signature: Optional[Dict[str, Any]] = None
    structure: Optional[Dict[str, Any]] = None
    entities: Optional[List[str]] = None
    key_terms: Optional[List[str]] = None


class MetadataExtractor:
    """Extract rich metadata from files"""
    
    def __init__(self, config: Optional[Dict[str, Any]] = None):
        self.config = config or {}
        self.extract_content_signature = self.config.get('extract_content_signature', True)
        self.extract_entities = self.config.get('extract_entities', True)
        self.extract_key_terms = self.config.get('extract_key_terms', True)
    
    def extract(self, file_path: str, file_id: Optional[str] = None) -> FileMetadata:
        """
        Extract comprehensive metadata from a file
        
        Args:
            file_path: Path to the file
            file_id: Optional unique ID (generated if not provided)
        
        Returns:
            FileMetadata object
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        stat = path.stat()
        file_id = file_id or self._generate_file_id(file_path)
        
        # Basic metadata
        metadata = FileMetadata(
            file_id=file_id,
            file_path=str(path.absolute()),
            file_name=path.name,
            file_type=self._detect_file_type(path.suffix),
            file_size=stat.st_size,
            extension=path.suffix.lower(),
            created_at=datetime.fromtimestamp(stat.st_ctime).isoformat(),
            modified_at=datetime.fromtimestamp(stat.st_mtime).isoformat()
        )
        
        # Extract format-specific metadata
        try:
            format_metadata = self._extract_format_metadata(file_path, path.suffix)
            if format_metadata:
                metadata.author = format_metadata.get('author')
                metadata.title = format_metadata.get('title')
                metadata.subject = format_metadata.get('subject')
                metadata.structure = format_metadata.get('structure')
        except Exception as e:
            # Continue even if format-specific extraction fails
            pass
        
        # Extract content signature
        if self.extract_content_signature:
            try:
                metadata.content_signature = self._extract_content_signature(file_path)
            except Exception:
                pass
        
        # Extract entities and key terms
        if self.extract_entities or self.extract_key_terms:
            try:
                content_data = self._extract_text_content(file_path)
                if content_data:
                    if self.extract_entities:
                        metadata.entities = self._extract_entities(content_data)
                    if self.extract_key_terms:
                        metadata.key_terms = self._extract_key_terms(content_data)
            except Exception:
                pass
        
        return metadata
    
    def _generate_file_id(self, file_path: str) -> str:
        """Generate unique file ID"""
        # Use file path + size + modified time for ID
        path = Path(file_path)
        stat = path.stat()
        content = f"{path.absolute()}{stat.st_size}{stat.st_mtime}"
        return hashlib.md5(content.encode()).hexdigest()
    
    def _detect_file_type(self, extension: str) -> str:
        """Detect file type from extension"""
        extension = extension.lower()
        type_map = {
            '.xlsx': 'excel',
            '.xls': 'excel',
            '.xlsm': 'excel',
            '.csv': 'csv',
            '.tsv': 'csv',
            '.json': 'json',
            '.pptx': 'powerpoint',
            '.ppt': 'powerpoint',
            '.docx': 'word',
            '.doc': 'word',
        }
        return type_map.get(extension, 'unknown')
    
    def _extract_format_metadata(self, file_path: str, extension: str) -> Optional[Dict[str, Any]]:
        """Extract format-specific metadata"""
        extension = extension.lower()
        
        if extension in ['.xlsx', '.xls', '.xlsm']:
            return self._extract_excel_metadata(file_path)
        elif extension in ['.pptx', '.ppt']:
            return self._extract_ppt_metadata(file_path)
        elif extension in ['.docx', '.doc']:
            return self._extract_word_metadata(file_path)
        
        return None
    
    def _extract_excel_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract Excel-specific metadata"""
        metadata = {}
        
        if not OPENPYXL_AVAILABLE:
            return metadata
        
        try:
            wb = load_workbook(file_path, read_only=True)
            
            # Document properties
            props = wb.properties
            if props:
                metadata['author'] = props.creator
                metadata['title'] = props.title
                metadata['subject'] = props.subject
            
            # Structure
            structure = {
                'sheet_count': len(wb.sheetnames),
                'sheet_names': wb.sheetnames
            }
            
            # Get column info from first sheet
            if wb.sheetnames:
                ws = wb[wb.sheetnames[0]]
                columns = []
                if ws.max_row > 0:
                    for cell in ws[1]:
                        if cell.value:
                            columns.append(str(cell.value))
                structure['columns'] = columns
                structure['row_count'] = ws.max_row
            
            metadata['structure'] = structure
            wb.close()
        except Exception:
            pass
        
        return metadata
    
    def _extract_ppt_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract PowerPoint-specific metadata"""
        metadata = {}
        
        if not PPTX_AVAILABLE:
            return metadata
        
        try:
            prs = Presentation(file_path)
            
            # Document properties
            core_props = prs.core_properties
            if core_props:
                metadata['author'] = core_props.author
                metadata['title'] = core_props.title
                metadata['subject'] = core_props.subject
            
            # Structure
            structure = {
                'slide_count': len(prs.slides),
                'slide_titles': []
            }
            
            for slide in prs.slides:
                # Try to get title from first shape
                title = None
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        title = shape.text.strip()
                        break
                structure['slide_titles'].append(title or f"Slide {len(structure['slide_titles']) + 1}")
            
            metadata['structure'] = structure
        except Exception:
            pass
        
        return metadata
    
    def _extract_word_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract Word-specific metadata"""
        metadata = {}
        
        if not DOCX_AVAILABLE:
            return metadata
        
        try:
            doc = Document(file_path)
            
            # Document properties
            core_props = doc.core_properties
            if core_props:
                metadata['author'] = core_props.author
                metadata['title'] = core_props.title
                metadata['subject'] = core_props.subject
            
            # Structure
            structure = {
                'paragraph_count': len(doc.paragraphs),
                'sections': []
            }
            
            # Extract section headers (paragraphs with heading style)
            for para in doc.paragraphs:
                if para.style.name.startswith('Heading'):
                    structure['sections'].append(para.text.strip())
            
            metadata['structure'] = structure
        except Exception:
            pass
        
        return metadata
    
    def _extract_content_signature(self, file_path: str) -> Dict[str, Any]:
        """Extract content signature (hash, key terms, etc.)"""
        try:
            # Read file content (first 1MB for hashing)
            with open(file_path, 'rb') as f:
                content = f.read(1024 * 1024)  # First 1MB
            
            content_hash = hashlib.sha256(content).hexdigest()
            
            return {
                'hash': content_hash,
                'hash_type': 'sha256'
            }
        except Exception:
            return {}
    
    def _extract_text_content(self, file_path: str) -> Optional[str]:
        """Extract text content from file for entity/key term extraction"""
        path = Path(file_path)
        extension = path.suffix.lower()
        
        try:
            if extension in ['.xlsx', '.xls', '.xlsm']:
                # Extract text from Excel
                df = pd.read_excel(file_path, nrows=100)  # First 100 rows
                return ' '.join(df.astype(str).values.flatten())
            elif extension == '.csv':
                df = pd.read_csv(file_path, nrows=100)
                return ' '.join(df.astype(str).values.flatten())
            elif extension == '.json':
                with open(file_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    return json.dumps(data)
            elif extension in ['.pptx', '.ppt']:
                if PPTX_AVAILABLE:
                    prs = Presentation(file_path)
                    texts = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, 'text'):
                                texts.append(shape.text)
                    return ' '.join(texts)
            elif extension in ['.docx', '.doc']:
                if DOCX_AVAILABLE:
                    doc = Document(file_path)
                    return ' '.join([para.text for para in doc.paragraphs])
        except Exception:
            pass
        
        return None
    
    def _extract_entities(self, text: str) -> List[str]:
        """Extract entities from text (simple pattern-based)"""
        entities = set()
        
        # Extract potential company names (Capitalized words)
        company_pattern = r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\b'
        matches = re.findall(company_pattern, text)
        entities.update([m for m in matches if len(m) > 3][:20])  # Top 20
        
        # Extract numbers that might be IDs
        id_pattern = r'\b[A-Z]{2,}\d{3,}\b'  # Like "PROJ1234"
        entities.update(re.findall(id_pattern, text)[:10])
        
        return sorted(list(entities))[:30]  # Return top 30
    
    def _extract_key_terms(self, text: str) -> List[str]:
        """Extract key terms from text"""
        # Simple frequency-based extraction
        words = re.findall(r'\b[a-zA-Z]{4,}\b', text.lower())
        
        # Common stop words to exclude
        stop_words = {'that', 'this', 'with', 'from', 'have', 'will', 'were', 'been', 'their', 'there'}
        
        # Count frequencies
        word_freq = {}
        for word in words:
            if word not in stop_words:
                word_freq[word] = word_freq.get(word, 0) + 1
        
        # Get top terms
        sorted_terms = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)
        return [term for term, freq in sorted_terms[:20]]  # Top 20 terms
    
    def to_dict(self, metadata: FileMetadata) -> Dict[str, Any]:
        """Convert FileMetadata to dictionary"""
        return asdict(metadata)

