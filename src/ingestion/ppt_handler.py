"""
PowerPoint file handler (.pptx)
"""
from typing import Dict, Any, List
from pathlib import Path
from pptx import Presentation
from .base_handler import BaseHandler


class PPTHandler(BaseHandler):
    """Handler for PowerPoint files"""
    
    def can_handle(self, file_path: str) -> bool:
        """Check if file is a PowerPoint file"""
        ext = Path(file_path).suffix.lower()
        return ext in ['.pptx', '.ppt']
    
    def get_supported_extensions(self) -> List[str]:
        return ['.pptx', '.ppt']
    
    def extract(self, source: str, **kwargs) -> Dict[str, Any]:
        """Extract data from PowerPoint file"""
        metadata = self.extract_metadata(source)
        
        try:
            prs = Presentation(source)
            
            slides_data = []
            all_text = []
            structured_data = []
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_text = []
                slide_data = {
                    'slide_number': slide_idx + 1,
                    'title': '',
                    'content': [],
                    'tables': [],
                    'images': 0
                }
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        slide_text.append(text)
                        
                        # Check if it's a title
                        if shape.shape_type == 1:  # Placeholder type
                            slide_data['title'] = text
                        else:
                            slide_data['content'].append(text)
                    
                    # Extract tables (check shape type - table type is 19 in python-pptx)
                    try:
                        from pptx.enum.shapes import MSO_SHAPE_TYPE
                        if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                            table_data = self._extract_table(shape.table)
                            slide_data['tables'].append(table_data)
                            structured_data.append({
                                'slide': slide_idx + 1,
                                'type': 'table',
                                'data': table_data
                            })
                    except (AttributeError, TypeError, ImportError):
                        # Shape doesn't have a table or can't import enum, skip it
                        pass
                    
                    # Count images
                    if hasattr(shape, "image"):
                        slide_data['images'] += 1
                
                slides_data.append(slide_data)
                all_text.append(f"Slide {slide_idx + 1}:\n" + "\n".join(slide_text))
            
            # Combine all text for LLM training
            full_text = "\n\n".join(all_text)
            
            metadata.update({
                'slide_count': len(prs.slides),
                'table_count': sum(len(s['tables']) for s in slides_data),
                'total_images': sum(s['images'] for s in slides_data)
            })
            
            return {
                'data': slides_data,
                'metadata': metadata,
                'text_content': [full_text],
                'structure': {
                    'format': 'powerpoint',
                    'slide_count': len(prs.slides),
                    'has_tables': any(s['tables'] for s in slides_data),
                    'structured_data': structured_data
                }
            }
            
        except Exception as e:
            raise ValueError(f"Error reading PowerPoint file {source}: {str(e)}")
    
    def _extract_table(self, table) -> List[Dict]:
        """Extract data from a PowerPoint table"""
        table_data = []
        headers = []
        
        # First row as headers
        if len(table.rows) > 0:
            headers = [cell.text.strip() for cell in table.rows[0].cells]
        
        # Extract rows
        for row_idx, row in enumerate(table.rows[1:], start=1):
            row_data = {}
            for col_idx, cell in enumerate(row.cells):
                header = headers[col_idx] if col_idx < len(headers) else f"Column_{col_idx}"
                row_data[header] = cell.text.strip()
            table_data.append(row_data)
        
        return table_data

